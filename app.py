# ===============================
# Universal Document RAG App
# ===============================

import os
import time
import streamlit as st
from dotenv import load_dotenv

from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd

from google import genai
from google.genai.errors import ClientError

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings


# ===============================
# Environment Setup (SAFE)
# ===============================

load_dotenv()

def get_api_key():
    """
    Priority:
    1. Streamlit Cloud Secrets
    2. Local .env file
    """
    try:
        return st.secrets["GOOGLE_API_KEY"]
    except Exception:
        return os.getenv("GOOGLE_API_KEY")

API_KEY = get_api_key()

if not API_KEY:
    st.error("❌ GOOGLE_API_KEY not found. Add it to Streamlit Secrets (cloud) or .env (local).")
    st.stop()


# ===============================
# Streamlit UI
# ===============================

st.set_page_config(
    page_title="Universal Document Chat (RAG)",
    layout="wide"
)

st.title("📂 Chat with Your Documents (RAG)")


# ===============================
# Session State Initialization
# ===============================

if "vectorstore" not in st.session_state:
    st.session_state.vectorstore = None

if "current_file" not in st.session_state:
    st.session_state.current_file = None

if "qa_cache" not in st.session_state:
    st.session_state.qa_cache = {}


# ===============================
# Helper Functions
# ===============================

def read_file(file):
    ext = file.name.split(".")[-1].lower()

    if ext == "pdf":
        reader = PdfReader(file)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if ext == "txt":
        return file.read().decode("utf-8")

    if ext == "docx":
        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext in ["xls", "xlsx"]:
        df = pd.read_excel(file)
        return df.to_string(index=False)

    if ext == "pptx":
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    return ""


def split_text(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=150
    )
    return splitter.split_text(text)


def create_vector_store(chunks):
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )
    return FAISS.from_texts(chunks, embeddings)


# ===============================
# File Upload
# ===============================

uploaded_file = st.file_uploader(
    "Upload a document (PDF, Word, Excel, Text, PPT)",
    type=["pdf", "txt", "docx", "xls", "xlsx", "pptx"]
)


# ===============================
# Document Processing
# ===============================

if uploaded_file:
    if st.session_state.current_file != uploaded_file.name:
        st.session_state.current_file = uploaded_file.name
        st.session_state.vectorstore = None
        st.session_state.qa_cache = {}

    if st.session_state.vectorstore is None:
        with st.spinner("Processing document..."):
            text = read_file(uploaded_file)

            if not text.strip():
                st.error("❌ No readable text found in the document.")
                st.stop()

            chunks = split_text(text)
            st.session_state.vectorstore = create_vector_store(chunks)

        st.success("✅ Document processed successfully")


# ===============================
# Question Answering (SAFE RAG)
# ===============================

with st.form("qa_form"):
    question = st.text_input("Ask a question from the document")
    submitted = st.form_submit_button("Ask")

if submitted and question and st.session_state.vectorstore:
    retriever = st.session_state.vectorstore.as_retriever(search_kwargs={"k": 3})
    docs = retriever.invoke(question)

    context = "\n\n".join(doc.page_content[:800] for doc in docs)

    cache_key = (question, context)

    if cache_key in st.session_state.qa_cache:
        st.subheader("📌 Answer (cached)")
        st.write(st.session_state.qa_cache[cache_key])
        st.stop()

    prompt = f"""
Answer the question using ONLY the context below.
If the answer is not present, say "Not found in the document".

Context:
{context}

Question:
{question}
"""

    try:
        client = genai.Client(api_key=API_KEY)

        response = client.models.generate_content(
            model="models/gemini-flash-latest",
            contents=prompt
        )

        st.session_state.qa_cache[cache_key] = response.text

        st.subheader("📌 Answer")
        st.write(response.text)

    except ClientError as e:
        if getattr(e, "code", None) == 429:
            st.warning("⚠️ API quota exceeded. Please wait and try again.")
            time.sleep(60)
        else:
            st.error(f"❌ Gemini API error: {e}")
